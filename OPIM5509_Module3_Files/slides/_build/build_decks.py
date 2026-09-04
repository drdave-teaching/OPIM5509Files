# -*- coding: utf-8 -*-
"""Branded OPIM 5509 M3.1 decks, built from code (python-pptx).
Theory: ConvNet Theory and Terminology. Math: one layer per slide, one big
calculation each, running total - 'calculations obvious, easy to digest'."""
from pptx import Presentation
from pptx.util import Inches as IN, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
import os, random, copy

GOLD  = RGBColor(0xF2, 0xB7, 0x05)
GOLDD = RGBColor(0x8C, 0x69, 0x08)
GOLDT = RGBColor(0x9A, 0x78, 0x0A)   # gold text on light
SOFT  = RGBColor(0xFA, 0xD8, 0x78)
DARK  = RGBColor(0x18, 0x1A, 0x20)
CARDDK= RGBColor(0x24, 0x26, 0x2E)
INK   = RGBColor(0x1C, 0x1E, 0x24)
GREY  = RGBColor(0x78, 0x7C, 0x86)
LGREY = RGBColor(0xA8, 0xAB, 0xB4)
PAPER = RGBColor(0xFC, 0xFC, 0xFA)
CARD  = RGBColor(0xFF, 0xFB, 0xEE)
CARDBRD=RGBColor(0xE3, 0xD9, 0xB8)
CHIPOFF=RGBColor(0xED, 0xEB, 0xE4)

MT = r'C:\Users\dww05002\Documents\DL_Fall2026_GeneralMaterials\M31_review\media_theory'
MM = r'C:\Users\dww05002\Documents\DL_Fall2026_GeneralMaterials\M31_review\media_math'
OUT = r'C:\Users\dww05002\Documents\DL_Fall2026_GeneralMaterials\M31_review\BRANDED'
os.makedirs(OUT, exist_ok=True)
W, H = 13.333, 7.5

def new_prs():
    p = Presentation(); p.slide_width = IN(W); p.slide_height = IN(H)
    return p

def blank(prs, dark=False):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, IN(W), IN(H))
    bg.fill.solid(); bg.fill.fore_color.rgb = DARK if dark else PAPER
    bg.line.fill.background(); bg.shadow.inherit = False
    return s

def rect(s, x, y, w, h, fill, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    r = s.shapes.add_shape(shape, IN(x), IN(y), IN(w), IN(h))
    if fill is None: r.fill.background()
    else: r.fill.solid(); r.fill.fore_color.rgb = fill
    if line is None: r.line.fill.background()
    else: r.line.color.rgb = line; r.line.width = Pt(lw)
    r.shadow.inherit = False
    return r

def text(s, x, y, w, h, runs, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, font='Arial', spacing=1.0):
    tb = s.shapes.add_textbox(IN(x), IN(y), IN(w), IN(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(runs, str): runs = [(runs, {})]
    if runs and isinstance(runs[0], tuple) and not isinstance(runs[0][1], dict):
        runs = [runs]
    # runs: list of paragraphs; each paragraph = list of (txt, overrides) or str
    first = True
    for para in (runs if isinstance(runs[0], list) else [runs]):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        if spacing != 1.0: p.line_spacing = spacing
        if isinstance(para, str): para = [(para, {})]
        for txt, ov in para:
            r = p.add_run(); r.text = txt
            r.font.name = ov.get('font', font)
            r.font.size = Pt(ov.get('size', size))
            r.font.bold = ov.get('bold', bold)
            r.font.color.rgb = ov.get('color', color)
            if ov.get('italic'): r.font.italic = True
    return tb

def pic(s, path, x, y, w=None, h=None):
    kw = {}
    if w: kw['width'] = IN(w)
    if h: kw['height'] = IN(h)
    return s.shapes.add_picture(path, IN(x), IN(y), **kw)

def pic_fit(s, path, x, y, maxw, maxh):
    """place centered inside box preserving aspect"""
    from PIL import Image
    iw, ih = Image.open(path).size
    sc = min(maxw/iw, maxh/ih)
    w, h = iw*sc, ih*sc
    return pic(s, path, x+(maxw-w)/2, y+(maxh-h)/2, w=w, h=h)

def baseline(s):
    rect(s, 0, H-0.14, W, 0.14, GOLD)
    rect(s, 0, H-0.18, W, 0.04, GOLDD)

def footer(s, n=None):
    baseline(s)
    text(s, 0.8, H-0.62, 6, 0.35, 'OPIM 5509 · Introduction to Deep Learning', size=13, color=GREY)
    if n: text(s, W-1.3, H-0.62, 0.6, 0.35, str(n), size=13, color=GREY, bold=True, align=PP_ALIGN.RIGHT)

def header(s, eyebrow, title, tsize=40):
    rect(s, 0.8, 0.55, 0.95, 0.055, GOLD)
    text(s, 0.8, 0.72, 11, 0.4, eyebrow.upper(), size=15, color=GOLDT, bold=True)
    text(s, 0.77, 1.05, 11.8, 1.0, title, size=tsize, color=INK, bold=True)

def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt

def nn_motif(s, x0=8.2, w=4.6):
    random.seed(5509)
    layers = [3, 5, 5, 1]
    lx = [x0 + w*f for f in (0.05, 0.38, 0.71, 0.98)]
    pos = [[(x, 1.0 + 5.5*(k+1)/(n+1)) for k in range(n)] for n, x in zip(layers, lx)]
    for li in range(len(pos)-1):
        for x1, y1 in pos[li]:
            for x2, y2 in pos[li+1]:
                c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, IN(x1), IN(y1), IN(x2), IN(y2))
                on = random.random() > 0.42
                c.line.color.rgb = GOLDD if on else RGBColor(0x3A,0x3D,0x46)
                c.line.width = Pt(1.4 if on else 0.9)
                c.shadow.inherit = False
    for li, layer in enumerate(pos):
        for x, y in layer:
            r = 0.13 if li < 3 else 0.22
            hot = random.random() > 0.3
            o = rect(s, x-r, y-r, 2*r, 2*r, GOLD if hot else RGBColor(0x34,0x37,0x40),
                     line=None if hot else RGBColor(0x5A,0x5D,0x68), lw=1.2, shape=MSO_SHAPE.OVAL)

def title_slide(prs, eyebrow, title_lines, sub):
    s = blank(prs, dark=True)
    nn_motif(s)
    rect(s, 0.85, 2.0, 1.4, 0.07, GOLD)
    text(s, 0.85, 2.25, 7, 0.5, eyebrow.upper(), size=20, color=GOLD, bold=True)
    text(s, 0.8, 2.75, 7.6, 2.6, [[ (ln, {}) ] for ln in title_lines], size=54, color=RGBColor(0xF5,0xF5,0xF7), bold=True, spacing=1.0)
    text(s, 0.85, 2.85+0.95*len(title_lines)+0.35, 7.2, 0.9, sub, size=17, color=LGREY)
    baseline(s)
    return s

def section_slide(prs, eyebrow, title, sub=''):
    s = blank(prs, dark=True)
    rect(s, 0.85, 2.6, 1.4, 0.07, GOLD)
    text(s, 0.85, 2.85, 10, 0.5, eyebrow.upper(), size=18, color=GOLD, bold=True)
    text(s, 0.8, 3.35, 11.5, 1.2, title, size=44, color=RGBColor(0xF5,0xF5,0xF7), bold=True)
    if sub: text(s, 0.85, 4.45, 11, 0.9, sub, size=18, color=LGREY)
    baseline(s)
    return s

def card(s, x, y, w, h, title=None, tcolor=GOLDD):
    c = rect(s, x, y, w, h, CARD, line=GOLD, lw=2.2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    c.adjustments[0] = 0.06
    if title: text(s, x+0.35, y+0.22, w-0.7, 0.4, title.upper(), size=16, color=tcolor, bold=True)
    return c

def bullets(s, x, y, w, items, size=16, gap=0.42, color=INK, dot=GOLD):
    for i, t in enumerate(items):
        rect(s, x, y+i*gap+0.09, 0.13, 0.13, dot, shape=MSO_SHAPE.OVAL)
        text(s, x+0.28, y+i*gap, w-0.3, gap, t, size=size, color=color)

def quote_card(s, x, y, w, h, lines, attribution):
    c = rect(s, x, y, w, h, DARK, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    c.adjustments[0] = 0.08
    text(s, x+0.35, y+0.28, w-0.7, h-0.9, [[(ln, {})] for ln in lines], size=19, color=SOFT, bold=True, spacing=1.05)
    text(s, x+0.35, y+h-0.55, w-0.7, 0.4, attribution, size=13, color=LGREY)

# =====================================================================
# THEORY DECK
# =====================================================================
prs = new_prs()

s = title_slide(prs, 'OPIM 5509  ·  Module 3.1', ['ConvNet Theory', 'and Terminology'],
                'Dr. Dave Wanik  ·  Introduction to Deep Learning  ·  University of Connecticut')
notes(s, 'Video 2-3 of M3.1. PowerPoint time - rare for this class, but ConvNet concepts deserve pictures.')

# 2. hook: pasta maker
s = blank(prs); header(s, 'Module 3.1 · ConvNet Theory', 'A ConvNet is a feature-making machine')
pic_fit(s, MT+r'\s02_img0.jpg', 7.3, 1.9, 5.2, 4.7)
bullets(s, 0.85, 2.35, 6.2, [
    'An image goes in... and gets convolved into FEATURE MAPS',
    'Each kernel highlights a different part of the image',
    'Useful regions "light up" - ears, whiskers, textures',
    'The model learns WHICH kernels help it predict'], size=17, gap=0.62)
quote_card(s, 0.85, 5.1, 6.0, 1.55, ['"Picture the conv layer as a pasta maker -', 'the dog comes out as feature maps."'], '- the hook from the 2022 lecture, still undefeated')
footer(s, 2); notes(s, 'The Steve Nouri GIF. Input dog -> pasta maker -> feature maps. Two kernels -> two maps, etc.')

# 3. images are data
s = blank(prs); header(s, 'Module 3.1 · ConvNet Theory', 'Images are data!')
pic_fit(s, MT+r'\s03_img2.png', 0.85, 2.0, 11.6, 4.0)
text(s, 0.85, 6.15, 11.6, 0.5, [[('B&W image = one 2D array of 0-255.   ', {'bold': True}), ('Color image = THREE stacked arrays (R, G, B channels). Pixels are just numbers - so it\'s all just math.', {})]], size=16, color=INK)
footer(s, 3); notes(s, 'Grayscale 0-255, single channel. RGB = 3 channels. This is why input_shape has that last dimension.')

# 4-6. kernels in action
for i, (name, img, extra) in enumerate([
    ('outline', 's04_img2.png', 'Constant kernel values - this one finds edges by amplifying differences between neighbors.'),
    ('emboss',  's05_img2.png', 'Gives the illusion of depth - emphasizes pixel differences in one direction.'),
    ('sharpen', 's06_img2.png', 'Amplifies the center pixel vs. its neighbors - the image "pops". Play live: setosa.io/ev/image-kernels')]):
    s = blank(prs); header(s, 'Module 3.1 · ConvNet Theory', f'Kernels in action: {name}')
    pic_fit(s, MT + '\\' + img, 0.85, 2.05, 11.6, 4.1)
    text(s, 0.85, 6.25, 11.6, 0.5, extra, size=16, color=GREY)
    footer(s, 4+i)
    notes(s, f'{name} kernel demo. In OUR ConvNets the kernel values are LEARNED, not fixed like Photoshop.')

# 7. convolution & stride
s = blank(prs); header(s, 'Module 3.1 · ConvNet Theory', 'Convolution: light downsampling')
pic_fit(s, MT+r'\s07_img1.gif', 0.85, 2.0, 6.2, 4.3)
card(s, 7.4, 2.0, 5.15, 2.5, 'the rules')
bullets(s, 7.75, 2.7, 4.6, ['Kernel slides with a stride of 1', 'Lots of OVERLAP between steps', 'Output = feature map, slightly smaller', 'New size: L - (M - 1)'], size=16, gap=0.44)
quote_card(s, 7.4, 4.75, 5.15, 1.55, ['"Convolution is a light', 'downsampling of information."'], '- say it with me')
footer(s, 7); notes(s, 'Green = image, yellow = kernel, element-wise multiply + sum. 32x32 with 3x3 -> 30x30.')

# 8. pooling
s = blank(prs); header(s, 'Module 3.1 · ConvNet Theory', 'Pooling: aggressive downsampling')
pic_fit(s, MT+r'\s08_img1.gif', 0.85, 2.0, 6.2, 4.3)
card(s, 7.4, 2.0, 5.15, 2.5, 'the rules')
bullets(s, 7.75, 2.7, 4.6, ['Stride = kernel size: NO overlap', 'Just takes the max in each window', 'New size: L / M  (2x2 kernel halves it)', 'ZERO trainable parameters'], size=16, gap=0.44)
quote_card(s, 7.4, 4.75, 5.15, 1.55, ['"It flips over like a domino -', 'no overlap, nothing learned."'], '- max pooling in one line')
footer(s, 8); notes(s, '4x4 -> 2x2 with a 2x2 kernel. Max or average - most people use max and move on.')

# 9. no more feature engineering
s = blank(prs); header(s, 'Module 3.1 · ConvNet Theory', 'No more feature engineering!')
pic_fit(s, MT+r'\s09_img2.png', 0.85, 2.0, 5.6, 4.4)
steps = [('INPUT DATA', 'raw pixels'), ('LOWER-LEVEL FEATURES', 'edges, curves, colors'),
         ('HIGHER-LEVEL FEATURES', 'ears, eyes, textures'), ('MODEL PREDICTION', 'cat or dog?')]
for i, (a, b) in enumerate(steps):
    y = 2.1 + i*1.02
    rect(s, 6.9, y, 0.32, 0.32, GOLD, shape=MSO_SHAPE.OVAL)
    text(s, 6.955, y+0.035, 0.25, 0.3, str(i+1), size=13, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    text(s, 7.4, y-0.05, 5.2, 0.45, a, size=17, color=INK, bold=True)
    text(s, 7.4, y+0.33, 5.2, 0.4, b, size=14, color=GREY)
text(s, 6.9, 6.25, 5.7, 0.8, 'Bonus: SPATIAL INVARIANCE - the cat can be anywhere in the frame. No more perfectly centered MNIST digits.', size=15, color=GOLDT, bold=True)
footer(s, 9); notes(s, 'The network builds its own features. Early layers = boring edges; deep layers = object parts.')

# 10. AlexNet
s = blank(prs); header(s, 'Module 3.1 · ConvNet Theory', 'AlexNet: a famous ConvNet')
pic_fit(s, MT+r'\s10_img2.jpg', 0.85, 2.2, 11.6, 3.9)
text(s, 0.85, 6.3, 11.6, 0.5, 'Won ImageNet 2012 and kicked off the deep learning era. Look closely - it is just conv, pool, conv, pool... then dense.', size=16, color=GREY)
footer(s, 10); notes(s, 'Point at the stages. Millions of images, 1000 classes.')

# 11. AlexNet annotated
s = blank(prs); header(s, 'Module 3.1 · ConvNet Theory', 'AlexNet = data prep + the NN you know')
pic_fit(s, MT+r'\s11_img2.jpg', 0.85, 2.3, 11.6, 3.7)
t1 = rect(s, 1.1, 6.1, 5.4, 0.62, CARD, line=GOLD, lw=2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, 1.3, 6.22, 5.1, 0.4, 'conv/pool stack = the DATA PREP we just learned', size=14, color=INK, bold=True)
t2 = rect(s, 7.0, 6.1, 5.4, 0.62, CARD, line=GOLD, lw=2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, 7.2, 6.22, 5.1, 0.4, 'FLATTEN -> dense layers = the NN from Module 2!', size=14, color=INK, bold=True)
footer(s, 11); notes(s, 'Demystify: a ConvNet is feature-making front end + the dense network they already built.')

# 12. cats vs dogs
s = blank(prs); header(s, 'Module 3.1 · ConvNet Theory', 'Coming attraction: Cats vs. Dogs')
pic_fit(s, MT+r'\s12_img2.gif', 0.85, 2.0, 7.6, 4.4)
card(s, 8.75, 2.0, 3.8, 3.0, 'module 3.2')
bullets(s, 9.05, 2.7, 3.3, ['The "hello world" of computer vision', 'Real, messy images', 'Data augmentation', 'Image generators'], size=15, gap=0.5)
footer(s, 12); notes(s, 'Teaser for the implementation week.')

# 13. channels
s = blank(prs); header(s, 'Module 3.1 · ConvNet Theory', 'A color image = 3 channels')
pic_fit(s, MT+r'\s13_img1.png', 0.85, 2.0, 6.0, 4.5)
text(s, 7.3, 2.6, 5.3, 2.6, [
  [('height 320 · width 400 · 3 channels', {'bold': True, 'size': 22})],
  [('', {})],
  [('The R, G and B arrays stack like pages in a book. A kernel on RGB input is really an M x N x 3 brick - one slice per channel, summed into ONE feature map.', {'size': 17})]], size=17, color=INK)
footer(s, 13); notes(s, 'Sets up the L in the parameter formula - and why channels confuse everyone in layer 2.')

# 14. filter size matters
s = blank(prs); header(s, 'Module 3.1 · ConvNet Theory', 'Filter size changes what you see')
labs = [('Original photo', 's14_img2.png'), ('3 x 3 filter: more detail, bigger map', 's14_img3.png'), ('15 x 15 filter: less detail, smaller map', 's14_img4.png')]
for i, (lab, img) in enumerate(labs):
    x = 0.85 + i*4.05
    pic_fit(s, MT + '\\' + img, x, 2.2, 3.7, 3.4)
    text(s, x, 5.75, 3.7, 0.7, lab, size=15, color=INK, bold=True, align=PP_ALIGN.CENTER)
footer(s, 14); notes(s, 'One layer of convolution with different kernel sizes. Small kernel = fine detail.')

# 15. pooling + relu
s = blank(prs); header(s, 'Module 3.1 · ConvNet Theory', 'Pooling shrinks the feature map')
pic_fit(s, MT+r'\s15_img2.png', 0.85, 2.2, 5.6, 3.9)
pic_fit(s, MT+r'\s15_img3.png', 6.85, 2.2, 5.6, 3.9)
text(s, 0.85, 6.2, 5.6, 0.5, 'Pooling layer (5 x 5)', size=15, color=INK, bold=True, align=PP_ALIGN.CENTER)
text(s, 6.85, 6.2, 5.6, 0.5, 'ReLU + pooling layer (5 x 5)', size=15, color=INK, bold=True, align=PP_ALIGN.CENTER)
text(s, 0.85, 6.7, 11.6, 0.4, 'You are welcome to put the ReLU activation right inside the Conv2D layer - check the docs.', size=13, color=GREY)
footer(s, 15); notes(s, '63x79x1 examples from the blog.')

# 16. parameters: conv2d
s = blank(prs); header(s, 'Module 3.1 · ConvNet Theory', 'Counting parameters: Conv2D')
card(s, 0.85, 2.1, 11.65, 1.5)
text(s, 1.2, 2.4, 11, 0.9, [[('Conv2D trainable parms  =  ((M · N · L) + B) · F', {'bold': True, 'size': 30, 'color': INK})]], size=30)
text(s, 1.2, 3.15, 11, 0.4, 'M, N = filter size   ·   L = channels coming INTO the layer   ·   B = 1 bias per filter   ·   F = feature maps created', size=14, color=GREY)
bullets(s, 0.85, 4.0, 11.5, [
    'Each filter is an M x N x L BRICK - one M x N slice of weights per input channel, summed into ONE output map',
    'One bias per filter - the kernel and its bias slide together (just like one bias per output node in dense layers)',
    'Example: L = 32 maps in, F = 64 maps out, 3 x 3 filter   ->   ((3 · 3 · 32) + 1) · 64  =  18,496 parameters'], size=17, gap=0.62)
quote_card(s, 0.85, 5.72, 11.65, 0.95, ['The L is the one that bites: deeper layers inherit L from the PREVIOUS layer\'s F.'], '')
footer(s, 16); notes(s, 'Slow down here. The brick picture is the whole game for layer-2 param counts.')

# 17. parameters: everything else
s = blank(prs); header(s, 'Module 3.1 · ConvNet Theory', 'Counting parameters: everything else')
cells = [('INPUT LAYER', '0 parameters', 'it just reads the image'),
         ('POOLING LAYER', '0 parameters', 'nothing is learned - it only shrinks the maps'),
         ('DENSE LAYER', '(n + 1) · m', 'n inputs, m outputs; the +1 is each output node\'s bias'),
         ('OUTPUT LAYER', '(n + 1) · m', 'it\'s just another dense layer')]
for i, (a, b, c) in enumerate(cells):
    x = 0.85 + (i % 2) * 6.0; y = 2.15 + (i // 2) * 2.2
    card(s, x, y, 5.65, 1.95, a)
    text(s, x+0.35, y+0.62, 5.0, 0.7, b, size=28, color=INK, bold=True)
    text(s, x+0.35, y+1.3, 5.0, 0.55, c, size=14, color=GREY)
footer(s, 17); notes(s, 'Pooling: zero. Dense: same formula as Module 2. No surprises.')

# 18. check it in code
s = blank(prs); header(s, 'Module 3.1 · ConvNet Theory', 'Check yourself against model.summary()')
pic_fit(s, MT+r'\s18_img3.png', 0.85, 2.4, 6.3, 1.2)
pic_fit(s, MT+r'\s18_img4.png', 0.85, 3.8, 6.3, 2.7)
card(s, 7.6, 2.4, 4.95, 3.4, 'the workflow')
bullets(s, 7.95, 3.1, 4.3, ['Read the code', 'Predict every output shape', 'Predict every param count', 'Run it - and check yourself', 'This is the exam question!'], size=16, gap=0.5)
footer(s, 18); notes(s, 'Bridge to the math deck + notebook where we do exactly this.')

# 19. what convnets learn
s = blank(prs); header(s, 'Module 3.1 · ConvNet Theory', 'What do ConvNets learn?')
pic_fit(s, MT+r'\s19_fill.jpg', 0.85, 2.0, 11.6, 4.3)
text(s, 0.85, 6.4, 11.6, 0.5, 'Early convolutions: edges and blobs.  Later convolutions: textures and object parts.  Then FLATTEN -> dense classifier.', size=16, color=GREY)
footer(s, 19); notes(s, 'Read the figure left to right: one input image -> early convs -> later convs -> classifier.')

# 20. section: advanced topics
s = section_slide(prs, 'A peek behind the curtain', 'Advanced Topics',
                  'For later lectures - but let\'s stay motivated. This is where ConvNets get exciting.')
notes(s, 'Teaser section - transfer learning comes in M3.3.')

# 21. keep the base
s = blank(prs); header(s, 'Advanced topics · Transfer learning', 'Keep the base, swap classifiers')
pic_fit(s, MT+r'\s22_fill.jpg', 0.85, 2.0, 11.6, 4.3)
text(s, 0.85, 6.4, 11.6, 0.5, 'Someone else trained the conv base on millions of images. You keep it, and train only your own classifier on top.', size=16, color=GREY)
footer(s, 21); notes(s, 'The big idea of transfer learning in one picture.')

# 22. fine tuning
s = blank(prs); header(s, 'Advanced topics · Transfer learning', 'Fine-tuning VGG16')
for i, img in enumerate(['s23_img2.png', 's23_img3.png', 's23_img4.png']):
    pic_fit(s, MT + '\\' + img, 0.85 + i*4.05, 2.1, 3.7, 4.3)
text(s, 0.85, 6.5, 11.6, 0.5, 'Unfreeze the last conv block and let backprop adapt those kernels to YOUR small dataset.', size=16, color=GREY)
footer(s, 22); notes(s, 'M3.3 preview.')

# 23. vgg16 kernels
s = blank(prs); header(s, 'Advanced topics · Transfer learning', 'Kernels from VGG16 - amazing!')
caps = ['Block 1: general textures, common to every problem', 'Block 2: patterns start to specialize', 'Block 4: textures specific to YOUR problem']
for i, (img, cap) in enumerate(zip(['s24_img2.png', 's24_img3.png', 's24_img4.png'], caps)):
    pic_fit(s, MT + '\\' + img, 0.85 + i*4.05, 2.1, 3.7, 3.7)
    text(s, 0.85 + i*4.05, 5.95, 3.7, 0.8, cap, size=14, color=INK, align=PP_ALIGN.CENTER)
footer(s, 23); notes(s, 'General -> specific as you go deeper. This is why transfer learning works.')

# 24. recap
s = blank(prs, dark=True)
rect(s, 0.85, 0.9, 1.4, 0.07, GOLD)
text(s, 0.85, 1.12, 10, 0.5, 'MODULE 3.1 · RECAP', size=18, color=GOLD, bold=True)
text(s, 0.8, 1.6, 11.5, 1.0, 'Five things to remember', size=44, color=RGBColor(0xF5,0xF5,0xF7), bold=True)
recap = ['Images are just numbers - B&W is one channel, RGB is three',
         'Convolution = LIGHT downsampling; kernels are LEARNED, and each filter is an M x N x L brick',
         'Pooling = AGGRESSIVE downsampling; no overlap, zero trainable parameters',
         'ConvNets build their own features: edges -> textures -> parts -> prediction (and they\'re spatially invariant)',
         'Conv2D parms = ((M · N · L) + B) · F  -  next up: use it to rebuild model.summary() by hand']
for i, t in enumerate(recap):
    y = 2.85 + i*0.82
    rect(s, 0.9, y+0.05, 0.34, 0.34, GOLD, shape=MSO_SHAPE.OVAL)
    text(s, 0.965, y+0.09, 0.26, 0.3, str(i+1), size=14, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    text(s, 1.5, y, 11.2, 0.75, t, size=17, color=RGBColor(0xE8,0xE8,0xEC))
baseline(s)
notes(s, 'Close the theory videos here; the math walkthrough continues in the next deck.')

prs.save(os.path.join(OUT, 'ConvNet Theory and Terminology.pptx'))
print('THEORY deck:', len(prs.slides.__iter__.__self__._sldIdLst), 'slides')
